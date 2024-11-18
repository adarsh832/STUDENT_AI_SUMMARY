import PyPDF2
from pptx import Presentation
import os

class FileHandler:
    def __init__(self):
        self.supported_extensions = {'.pdf', '.pptx', '.ppt'}
        
    def read_file(self, file_path: str) -> str:
        """Read and extract text from supported files"""
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
            
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {file_ext}")
            
        if file_ext == '.pdf':
            return self._read_pdf(file_path)
        elif file_ext in ['.pptx', '.ppt']:
            return self._read_pptx(file_path)
            
    def _read_pdf(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text = []
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text.append(page.extract_text())
        return '\n'.join(text)
        
    def _read_pptx(self, file_path: str) -> str:
        """Extract text from PowerPoint file"""
        text = []
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text) 